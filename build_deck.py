#!/usr/bin/env python3
"""
デッキ組み立てパイプライン
================================================================
deck spec(JSON) を読み込み、型ごとに背景プールから選択し、
slide_render.py の9型レンダラーで描画、PPTXにまとめる。

実行時に画像生成APIは呼ばない（背景はプールから「選ぶ」だけ）。
= 即時・追加コストゼロ・世界観は事前検品済み。

使い方:
  python build_deck.py specs/deck_taniuchi.json
  python build_deck.py specs/deck_taniuchi.json --color terracotta
"""
import argparse, json, os
from slide_render import (
    render_cover, render_list, render_statement, render_message_over_image,
    render_multi_image, render_before_after, render_three_column,
    render_person_text, render_cta, render_content, render_offer,
    render_process, render_choice,
)
from pptx import Presentation
from pptx.util import Inches

ROOT_DIR = os.path.dirname(__file__)
MESSAGE_FALLBACKS = [
    os.path.join(ROOT_DIR, "assets", "message1.webp"),
    os.path.join(ROOT_DIR, "assets", "message2.webp"),
]

# 型名 → レンダラー関数
RENDERERS = {
    "cover":       render_cover,
    "list":        render_list,
    "statement":   render_statement,
    "message":     render_message_over_image,
    "multi_image": render_multi_image,
    "before_after": render_before_after,
    "three_column": render_three_column,
    "content":     render_content,
    "process":     render_process,
    "choice":      render_choice,
    "offer":       render_offer,
    "person":      render_person_text,
    "cta":         render_cta,
}

BACKGROUND_ALIASES = {
    "before_after": "multi_image",
    "three_column": "statement",
    "content": "multi_image",
    "process": "statement",
    "choice": "statement",
    "offer": "cta",
}

BACKGROUND_STYLES = {
    "open_center",
    "paper_stage",
    "postcard_cta",
    "photo_frame",
}

DEFAULT_BACKGROUND_STYLE = {
    "cover": "paper_stage",
    "list": "paper_stage",
    "statement": "paper_stage",
    "message": "paper_stage",
    "process": "open_center",
    "choice": "open_center",
    "content": "open_center",
    "offer": "open_center",
    "cta": "postcard_cta",
    "multi_image": "photo_frame",
    "before_after": "photo_frame",
    "three_column": "open_center",
    "person": "photo_frame",
}

# message スライド用の一時背景（assets 内の webp）
MESSAGE_FALLBACKS = [
    "assets/message1.webp",
    "assets/message2.webp",
]


class BackgroundPool:
    """背景プール。新style指定と旧色×型指定の両方を扱う。"""
    def __init__(self, pool_dir):
        self.pool, self.style_pool, self.idx = {}, {}, {}
        self.pool_dir = pool_dir
        known_types = sorted(
            set(RENDERERS) | set(BACKGROUND_ALIASES.values()),
            key=len,
            reverse=True,
        )
        for color in os.listdir(pool_dir):
            color_path = os.path.join(pool_dir, color)
            if not os.path.isdir(color_path):
                continue
            if color in BACKGROUND_STYLES:
                self._load_style_dir(color, color_path)
                continue
            for f in sorted(os.listdir(color_path)):
                if not f.lower().endswith((".png", ".jpg", ".jpeg")):
                    continue
                stem = os.path.splitext(f)[0]
                slide_type = next(
                    (t for t in known_types if stem == t or stem.startswith(f"{t}_")),
                    None,
                )
                if not slide_type:
                    continue
                self.pool.setdefault(color, {}).setdefault(slide_type, []).append(os.path.join(color_path, f))
        for color, by_type in self.pool.items():
            for slide_type, imgs in by_type.items():
                imgs.sort()
                self.idx[(color, slide_type)] = 0
        for style, by_color in self.style_pool.items():
            for color, by_tone in by_color.items():
                for tone, imgs in by_tone.items():
                    imgs.sort()
                    self.idx[(style, color, tone)] = 0

    def _load_style_dir(self, style, style_path):
        for f in sorted(os.listdir(style_path)):
            if not f.lower().endswith((".png", ".jpg", ".jpeg")):
                continue
            stem = os.path.splitext(f)[0]
            tone = None
            for suffix in ("_normal", "_reverse"):
                if stem.endswith(suffix):
                    color = stem[:-len(suffix)]
                    tone = suffix[1:]
                    break
            if not tone or not color:
                continue
            self.style_pool.setdefault(style, {}).setdefault(color, {}).setdefault(tone, []).append(
                os.path.join(style_path, f)
            )

    def has(self, slide_type, color):
        return (
            (color in self.pool and slide_type in self.pool[color])
            or ("default" in self.pool and slide_type in self.pool["default"])
        )

    def pick(self, slide_type, color):
        if not self.has(slide_type, color):
            raise ValueError(
                f"背景プールに色 '{color}'・型 '{slide_type}' がありません"
                f"（{self.pool_dir}/{color}/{slide_type}.png を確認）"
            )
        actual_color = color if color in self.pool and slide_type in self.pool[color] else "default"
        imgs = self.pool[actual_color][slide_type]
        key = (actual_color, slide_type)
        img = imgs[self.idx[key] % len(imgs)]
        self.idx[key] += 1
        return img

    def has_style(self, style, color, tone="normal"):
        return (
            style in self.style_pool
            and (
                (color in self.style_pool[style] and tone in self.style_pool[style][color])
                or ("default" in self.style_pool[style] and tone in self.style_pool[style]["default"])
            )
        )

    def pick_style(self, style, color, tone="normal"):
        if not self.has_style(style, color, tone):
            raise ValueError(
                f"背景プールにstyle '{style}'・色 '{color}'・tone '{tone}' がありません"
                f"（{self.pool_dir}/{style}/{color}_{tone}.png を確認）"
            )
        actual_color = (
            color
            if color in self.style_pool[style] and tone in self.style_pool[style][color]
            else "default"
        )
        imgs = self.style_pool[style][actual_color][tone]
        key = (style, actual_color, tone)
        img = imgs[self.idx[key] % len(imgs)]
        self.idx[key] += 1
        return img


def resolve_bg_type(stype, pool, color):
    if pool.has(stype, color):
        return stype
    return BACKGROUND_ALIASES.get(stype, stype)

def resolve_background(slide, stype, pool, default_color):
    bg_spec = slide.get("background")
    slide_color = slide.get("color", default_color)

    if isinstance(bg_spec, str):
        return bg_spec, slide_color, bg_spec

    if isinstance(bg_spec, dict):
        style = bg_spec.get("style") or DEFAULT_BACKGROUND_STYLE.get(stype)
        bg_color = bg_spec.get("color") or slide_color
        tone = bg_spec.get("tone") or "normal"
        if style and pool.has_style(style, bg_color, tone):
            bg = pool.pick_style(style, bg_color, tone)
            return bg, bg_color, f"{style}/{bg_color}_{tone}"
        if style:
            print(
                f"  warn: background style not found: style={style}, color={bg_color}, tone={tone}; "
                "legacy backgroundへフォールバック"
            )

    style = DEFAULT_BACKGROUND_STYLE.get(stype)
    if style and pool.has_style(style, slide_color, "normal"):
        bg = pool.pick_style(style, slide_color, "normal")
        return bg, slide_color, f"{style}/{slide_color}_normal"
    if style and pool.has_style(style, slide_color, "reverse"):
        bg = pool.pick_style(style, slide_color, "reverse")
        return bg, slide_color, f"{style}/{slide_color}_reverse"

    bg_type = resolve_bg_type(stype, pool, slide_color)
    bg = pool.pick(bg_type, slide_color)
    return bg, slide_color, f"{slide_color}/{bg_type}"

def is_reverse_background(bg_label):
    return (
        isinstance(bg_label, str)
        and (
            bg_label.endswith("_reverse")
            or "_reverse." in bg_label
        )
    )

def text_lines_from_content(content, fallback="Generated slide"):
    if not isinstance(content, dict):
        return [str(content or fallback)]
    for key in ("lines", "title", "heading", "body", "subtext"):
        value = content.get(key)
        if isinstance(value, list):
            return [str(v) for v in value if str(v).strip()] or [fallback]
        if isinstance(value, str) and value.strip():
            return [value]
    items = content.get("items") or []
    lines = []
    for item in items:
        if isinstance(item, dict):
            title = item.get("title") or item.get("label") or item.get("name")
            body = item.get("body") or item.get("text") or item.get("caption")
            line = " / ".join(str(v) for v in (title, body) if v)
            if line:
                lines.append(line)
        elif str(item).strip():
            lines.append(str(item))
    return lines or [fallback]


def items_to_steps(items):
    steps = []
    for idx, item in enumerate(items or [], 1):
        if isinstance(item, dict):
            title = item.get("title") or item.get("label") or item.get("name") or f"Step {idx}"
            body = item.get("body") or item.get("text") or item.get("caption") or ""
        else:
            title = f"Step {idx}"
            body = str(item)
        steps.append({"label": str(idx), "title": str(title), "body": str(body)})
    return steps


def normalize_slide(stype, content):
    if not isinstance(content, dict):
        return "statement", {"lines": [str(content or "Generated slide")]}

    if stype == "content":
        items = content.get("items") or []
        has_image_items = any(isinstance(item, dict) and item.get("image") for item in items)
        if items and not has_image_items:
            return "process", {
                "title": content.get("title") or content.get("heading") or "Overview",
                "steps": items_to_steps(items),
                "columns": min(max(len(items), 1), 4),
            }
        if not items:
            return "statement", {"lines": text_lines_from_content(content, "Overview")}

    if stype == "process" and not content.get("steps"):
        items = content.get("items") or []
        if items:
            content = dict(content)
            content["steps"] = items_to_steps(items)

    if stype == "statement" and not content.get("lines"):
        content = dict(content)
        content["lines"] = text_lines_from_content(content)

    if stype == "cta":
        content = dict(content)
        content.setdefault("heading", content.get("title") or "Next Step")
        content.setdefault("qr_data", "https://gen-deck.onrender.com")

    return stype, content


def build_deck(spec_path, pool_dir="backgrounds", output_pptx=None, work_dir="_slides", color="terracotta"):
    os.makedirs(work_dir, exist_ok=True)
    with open(spec_path, encoding="utf-8") as f:
        spec = json.load(f)

    title = spec.get("deck_title", "deck")
    pool = BackgroundPool(pool_dir)
    default_color = spec.get("color", color)
    output_pptx = output_pptx or f"{title}_{default_color}.pptx"
    slides = spec["slides"]
    print(f"=== 「{title}」 {len(slides)}枚を組み立て（default color: {default_color}）===")

    pngs = []
    message_idx = 0
    valid_message_fallbacks = [p for p in MESSAGE_FALLBACKS if os.path.exists(p)]
    for i, sl in enumerate(slides, 1):
        stype = sl["type"]
        stype, content = normalize_slide(stype, sl["content"])
        slide_color = sl.get("color", default_color)
        if stype not in RENDERERS:
            stype = "statement"
            content = {"lines": text_lines_from_content(content, f"Slide {i}")}
        # 背景：明示指定があればそれ、なければプールから選択
        if stype == "message":
            if isinstance(sl.get("background"), str):
                bg = sl["background"]
            elif pool.has("message", slide_color):
                bg = pool.pick("message", slide_color)
            elif valid_message_fallbacks:
                bg = valid_message_fallbacks[message_idx % len(valid_message_fallbacks)]
                message_idx += 1
            else:
                bg = pool.pick("statement", slide_color)
            bg_label = bg
        else:
            bg, slide_color, bg_label = resolve_background(sl, stype, pool, default_color)
        if is_reverse_background(bg_label) and isinstance(content, dict) and "reverse" not in content:
            content = dict(content)
            content["reverse"] = True
        out = os.path.join(work_dir, f"slide_{i:02d}.png")
        print(f"  slide {i:02d}: type={stype}, color={slide_color}, bg={bg_label}")
        RENDERERS[stype](bg, out, content)
        pngs.append(out)

    # --- PPTX化（16:9・各スライドにフルブリードで画像を配置）---
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]   # 空白レイアウト
    for png in pngs:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(png, 0, 0,
                             width=prs.slide_width, height=prs.slide_height)
    prs.save(output_pptx)
    print(f"=== 完了: {output_pptx}（{len(pngs)}枚）===")
    return output_pptx


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("spec", nargs="?", default="specs/deck_taniuchi.json")
    ap.add_argument("--color", default="terracotta",
                    help="背景色プール。terracotta / cream / deep_green など")
    ap.add_argument("--pool-dir", default="backgrounds")
    ap.add_argument("--output", default=None)
    ap.add_argument("--work-dir", default="_slides")
    args = ap.parse_args()
    build_deck(args.spec, pool_dir=args.pool_dir, output_pptx=args.output,
               work_dir=args.work_dir, color=args.color)
